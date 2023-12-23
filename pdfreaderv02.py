import streamlit as st
import pandas as pd
import nltk
import re
import os
import seaborn as sns
import matplotlib.pyplot as plt
from docx import Document
from pptx import Presentation
from pypdf import PdfReader
from PIL import Image
from wordcloud import WordCloud
from transformers import pipeline

import warnings
warnings.filterwarnings("ignore")
st.set_option('deprecation.showPyplotGlobalUse', False)

##########################################################

# Function to extract text from a Word document
def extract_text_from_docx(file):
    doc = Document(file)
    text = ''
    for paragraph in doc.paragraphs:
        text += paragraph.text + '\n'
    return text

##########################################################

# Function to extract text from a PowerPoint file
def extract_text_from_pptx(file):
    prs = Presentation(file)
    text = ''
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text += shape.text + '\n'
    return text

##########################################################

# Function to extract text from a PDF file
def extract_text_from_pdf(file):
    pdf_reader = PdfReader(file) 
    text = ''
    num_pages = len(pdf_reader.pages)
    for page in range(num_pages):
        text += pdf_reader.pages[page].extract_text()
    return text

###########################################################

# Function to determine file extension and choose extraction method
def choose_extraction(file):
    file_name, file_extension = os.path.splitext(file.name)
    if file_extension.lower() == '.pdf':
        return extract_text_from_pdf(file)
    elif file_extension.lower() == '.docx':
        return extract_text_from_docx(file)
    else:
        return "Unsupported file format"
    
###########################################################

# Updated function to determine file extension and choose extraction method
def choose_extraction(file):
    file_name, file_extension = os.path.splitext(file.name)
    if file_extension.lower() == '.pdf':
        return extract_text_from_pdf(file)
    elif file_extension.lower() == '.docx':
        return extract_text_from_docx(file)
    elif file_extension.lower() == '.pptx':
        return extract_text_from_pptx(file)
    else:
        return "Unsupported file format"

############################################################
nltk.download("stopwords")
############################################################

model_name = "deepset/roberta-base-squad2"

############################################################
#Function to calculate word frequencies 
def generate_word_freq(text):
    token=re.findall('\w+', text)
    token = [c for c in token if not c.isnumeric()] #remove numeric values
    words=[]
    for word in token:
        words.append(word.lower())
    sw=nltk.corpus.stopwords.words('english')
    
    # get the list without stop words
    words_ne=[]
    for word in words:
        if word not in sw:
            words_ne.append(word)
    sns.set_style('darkgrid')
    nlp_words=nltk.FreqDist(words_ne)
    nlp_words.plot(15);

############################################################
#Function for word cloud 
def generate_wordcloud(text):
    wordcloud = WordCloud(max_font_size=30, max_words=100, random_state=123)
    wordcloud.generate(text)
    plt.imshow(wordcloud, interpolation='bilinear')
    plt.axis("off")
    plt.show()

############################################################
# Function to generate a summary using Hugging Face Transformers
def generate_summary(text):
    summarization_pipeline = pipeline("summarization") 
    summary = summarization_pipeline(text, max_length=50, min_length=30, do_sample=True, clean_up_tokenization_spaces=True)
    return summary[0]['summary_text']

# Function to classify text using Hugging Face Transformers
def generate_sentiment(text):
    classifier = pipeline("text-classification")
    outputs = classifier(text)
    sent_df = pd.DataFrame(outputs)
    return sent_df

# Function to answer questions using Hugging Face Transformers
def answer_questions(text, question):
    question_answering_pipeline = pipeline("question-answering", model=model_name, tokenizer=model_name)
    answer = question_answering_pipeline(question=question, context=text)
    return answer['answer']

#############################################################
#Streamlit app
def main():
    st.title("Text Summariztion, Question Answering App")
    st.header("Developed by: Al Yazdani")
    uploaded_file = st.file_uploader("Upload a file (pdf, word, powerpoint)", type=["pdf","docx","pptx"])

    if uploaded_file is not None:
        st.subheader("Sample Contents:")
        extracted_text = choose_extraction(uploaded_file)
        sample_text = extracted_text[0:100]
        st.text(sample_text)
       
        st.subheader("Word Cloud")
        #wordcloud = WordCloud(max_font_size=30, max_words=100, random_state=123)
        generate_wordcloud(extracted_text)
        st.pyplot()

        st.subheader("Most frequent words")
        st.pyplot(generate_word_freq(extracted_text))

        st.subheader("Sentiment:")
        sentiment = generate_sentiment(extracted_text)
        st.write(sentiment)

        st.subheader("Highlight:")
        summary = generate_summary(extracted_text)
        st.write(summary)

        st.subheader("Question Answering:")
        question = st.text_input("Ask a question about the document:")
        if st.button("Get Answer"):
            if question:
                answer = answer_questions(extracted_text, question)
                st.write("Answer:", answer)
            else:
                st.write("Please enter a question.")

if __name__ == "__main__":
    main()
